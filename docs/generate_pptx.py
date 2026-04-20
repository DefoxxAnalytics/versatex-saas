"""
Generate professionally formatted PowerPoint slide deck from Management Introduction
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

NAVY = RGBColor(0x1e, 0x3a, 0x5f)
TEAL = RGBColor(0x0d, 0x94, 0x88)
LIGHT_GRAY = RGBColor(0xf8, 0xfa, 0xfc)
DARK_GRAY = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xff, 0xff, 0xff)


def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_shape_fill(bg, NAVY)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), Inches(13.333), Inches(0.1))
    set_shape_fill(accent, TEAL)
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.paragraphs[0].text = subtitle
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.color.rgb = TEAL
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf3 = footer_box.text_frame
    tf3.paragraphs[0].text = f"Versatex / Defoxx Analytics  |  {datetime.now().strftime('%B %Y')}"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_num, section_title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    set_shape_fill(header, NAVY)
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = header_text.text_frame
    tf.paragraphs[0].text = "Versatex Analytics"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf2 = num_box.text_frame
    tf2.paragraphs[0].text = f"0{section_num}" if section_num < 10 else str(section_num)
    tf2.paragraphs[0].font.size = Pt(72)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = TEAL
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
    tf3 = title_box.text_frame
    tf3.paragraphs[0].text = section_title
    tf3.paragraphs[0].font.size = Pt(36)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = NAVY
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, highlight_box=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    set_shape_fill(header, NAVY)
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = header_text.text_frame
    tf.paragraphs[0].text = "Versatex Analytics"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.7))
    tf2 = title_box.text_frame
    tf2.paragraphs[0].text = title
    tf2.paragraphs[0].font.size = Pt(28)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.15), Inches(2), Inches(0.06))
    set_shape_fill(accent, TEAL)
    accent.line.fill.background()

    if highlight_box:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(2.5), Inches(4.8), Inches(1.5))
        set_shape_fill(box, LIGHT_GRAY)
        box.line.color.rgb = TEAL
        box.line.width = Pt(2)

        box_text = slide.shapes.add_textbox(Inches(8.2), Inches(2.7), Inches(4.4), Inches(1.3))
        tf_box = box_text.text_frame
        tf_box.word_wrap = True
        p = tf_box.paragraphs[0]
        p.text = "Business Value"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEAL
        p2 = tf_box.add_paragraph()
        p2.text = highlight_box
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_GRAY
        p2.space_before = Pt(6)

        bullet_width = 7
    else:
        bullet_width = 12

    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(bullet_width), Inches(4.5))
    tf3 = bullet_box.text_frame
    tf3.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    return slide


def add_module_slide(prs, module_name, purpose, features, business_value):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    set_shape_fill(header, NAVY)
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = header_text.text_frame
    tf.paragraphs[0].text = "Versatex Analytics"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.7))
    tf2 = title_box.text_frame
    tf2.paragraphs[0].text = module_name
    tf2.paragraphs[0].font.size = Pt(32)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = TEAL

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(2), Inches(0.06))
    set_shape_fill(accent, TEAL)
    accent.line.fill.background()

    purpose_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7.5), Inches(0.5))
    tf3 = purpose_box.text_frame
    tf3.paragraphs[0].text = purpose
    tf3.paragraphs[0].font.size = Pt(16)
    tf3.paragraphs[0].font.italic = True
    tf3.paragraphs[0].font.color.rgb = DARK_GRAY

    features_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(7.5), Inches(3.5))
    tf4 = features_box.text_frame
    tf4.word_wrap = True

    for i, feature in enumerate(features):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    value_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.5), Inches(4.3), Inches(3.5))
    set_shape_fill(value_box, RGBColor(0xe8, 0xf5, 0xf3))
    value_box.line.color.rgb = TEAL
    value_box.line.width = Pt(2)

    value_text = slide.shapes.add_textbox(Inches(8.7), Inches(2.7), Inches(3.9), Inches(3.1))
    tf5 = value_text.text_frame
    tf5.word_wrap = True
    p1 = tf5.paragraphs[0]
    p1.text = "Business Value"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = TEAL
    p2 = tf5.add_paragraph()
    p2.text = business_value
    p2.font.size = Pt(15)
    p2.font.color.rgb = NAVY
    p2.space_before = Pt(12)

    return slide


def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    set_shape_fill(header, NAVY)
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = header_text.text_frame
    tf.paragraphs[0].text = "Versatex Analytics"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.7))
    tf2 = title_box.text_frame
    tf2.paragraphs[0].text = title
    tf2.paragraphs[0].font.size = Pt(28)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.15), Inches(2), Inches(0.06))
    set_shape_fill(accent, TEAL)
    accent.line.fill.background()

    num_cols = len(headers)
    num_rows = len(rows) + 1
    col_width = 12 / num_cols

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(2.5), Inches(12), Inches(0.5 * num_rows)).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            if col_idx == 0:
                p.font.bold = True

    return slide


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "Versatex Analytics", "Enterprise Procurement Analytics Platform")

    add_content_slide(prs, "Executive Summary", [
        "Enterprise-grade procurement analytics platform",
        "Transform raw procurement data into actionable insights",
        "Comprehensive visibility into organizational spending",
        "End-to-end procure-to-pay process analytics"
    ])

    add_content_slide(prs, "Key Value Propositions", [
        "Centralized spend visibility across all categories and suppliers",
        "AI-powered recommendations with measurable ROI tracking",
        "End-to-end procure-to-pay process analytics",
        "Automated reporting with customizable scheduling"
    ])

    add_section_slide(prs, 1, "Core Procurement Dashboard")

    add_module_slide(prs, "Overview Dashboard",
        "Executive-level snapshot of procurement health",
        [
            "Real-time KPI cards (Total Spend, Supplier Count, Transactions)",
            "Interactive charts: Category, Supplier, Monthly Trends",
            "Drill-down capability for detailed breakdowns",
            "Date range filtering with quick presets"
        ],
        "Single-pane view for executives to monitor procurement performance at a glance"
    )

    add_module_slide(prs, "Categories Module",
        "Analyze spending by procurement category",
        [
            "Category hierarchy visualization",
            "Subcategory breakdown analysis",
            "Supplier distribution per category",
            "Risk scoring and concentration analysis"
        ],
        "Identify category-level savings opportunities and consolidation targets"
    )

    add_module_slide(prs, "Suppliers Module",
        "Supplier portfolio management and risk assessment",
        [
            "HHI concentration risk indicators",
            "Supplier ranking by spend volume",
            "Performance metrics and scorecards",
            "Search and filtering capabilities"
        ],
        "Mitigate supplier risk, identify strategic vs. tactical suppliers"
    )

    add_section_slide(prs, 2, "Advanced Spend Analytics")

    add_module_slide(prs, "Pareto Analysis",
        "Apply 80/20 rule to supplier base",
        [
            "Cumulative spend visualization",
            "Supplier classification: Critical (80%), Important (15%), Tail (5%)",
            "Drill-down to individual supplier details"
        ],
        "Focus negotiation efforts on high-impact suppliers"
    )

    add_module_slide(prs, "Spend Stratification",
        "Kraljic matrix-style spend segmentation",
        [
            "Four segments: Strategic, Leverage, Routine, Tactical",
            "Spend band analysis ($0-10K to $100K+)",
            "Category and supplier drill-downs"
        ],
        "Develop category-specific procurement strategies"
    )

    add_module_slide(prs, "Seasonality Analysis",
        "Identify spending patterns across time",
        [
            "Monthly spending heatmap",
            "Seasonal index calculations",
            "Peak/trough identification",
            "Fiscal vs. calendar year toggle"
        ],
        "Optimize procurement timing and budget planning"
    )

    add_module_slide(prs, "Year-over-Year Comparison",
        "Compare spending trends across years",
        [
            "Dual-year variance analysis",
            "Top gainers and decliners identification",
            "Category and supplier drill-downs"
        ],
        "Track procurement efficiency improvements over time"
    )

    add_module_slide(prs, "Tail Spend Analysis",
        "Identify fragmented, low-value spending",
        [
            "Adjustable threshold slider for classification",
            "Vendor fragmentation metrics",
            "Consolidation opportunity identification"
        ],
        "Reduce administrative overhead, consolidate vendor base"
    )

    add_section_slide(prs, 3, "AI & Predictive Analytics")

    add_module_slide(prs, "AI Insights",
        "AI-powered recommendations and ROI tracking",
        [
            "Cost optimization recommendations",
            "Supplier risk alerts and anomaly detection",
            "Consolidation opportunities",
            "ROI tracking: measure actual vs. projected savings"
        ],
        "Data-driven decision support with measurable outcomes"
    )

    add_module_slide(prs, "Predictive Analytics",
        "Forecast future spending patterns",
        [
            "Spend prediction models",
            "Demand forecasting",
            "Model accuracy metrics (MAPE, R²)",
            "Confidence intervals"
        ],
        "Improve budgeting accuracy and cash flow planning"
    )

    add_section_slide(prs, 4, "Risk & Compliance")

    add_module_slide(prs, "Contract Optimization",
        "Maximize contract value and coverage",
        [
            "Contract portfolio overview",
            "Utilization tracking (% of contracted spend)",
            "Expiration alerts",
            "Category breakdown per contract"
        ],
        "Reduce off-contract spending, improve contract leverage"
    )

    add_module_slide(prs, "Maverick Spend",
        "Identify and resolve policy violations",
        [
            "Policy violation listing",
            "Batch resolution capability",
            "Notes and documentation tracking",
            "Trend analysis"
        ],
        "Enforce procurement policies, reduce compliance risk"
    )

    add_section_slide(prs, 5, "Procure-to-Pay Analytics")

    add_content_slide(prs, "P2P Analytics Overview", [
        "P2P Cycle Analysis — End-to-end process timing and bottlenecks",
        "3-Way Matching — PO vs. GR vs. Invoice with exception management",
        "Invoice Aging — AP aging buckets and DPO trends",
        "Requisitions — PR workflow and approval analysis",
        "Purchase Orders — PO compliance and amendment tracking",
        "Supplier Payments — Payment performance scorecards"
    ])

    add_module_slide(prs, "P2P Cycle Analysis",
        "End-to-end process visibility",
        [
            "Full cycle timing: PR → PO → GR → Invoice → Payment",
            "Bottleneck identification",
            "Process funnel visualization",
            "Category and supplier trends"
        ],
        "Reduce cycle times, improve operational efficiency"
    )

    add_module_slide(prs, "3-Way Matching",
        "Invoice accuracy and exception management",
        [
            "PO vs. GR vs. Invoice matching",
            "Exception dashboard",
            "Price and quantity variance analysis",
            "Individual and bulk resolution"
        ],
        "Prevent overpayments, reduce invoice disputes"
    )

    add_module_slide(prs, "Invoice Aging",
        "Accounts payable management",
        [
            "Aging buckets: Current, 31-60, 61-90, 90+ days",
            "Days Payable Outstanding (DPO) trends",
            "Cash flow forecasting",
            "Supplier aging breakdown"
        ],
        "Optimize working capital, maintain supplier relationships"
    )

    add_section_slide(prs, 6, "Reporting & Administration")

    add_content_slide(prs, "Reports Module", [
        "14+ Report Types: Executive Summary, Spend Analysis, Pareto, and more",
        "P2P Reports: PR Status, PO Compliance, AP Aging",
        "Output Formats: PDF, Excel, CSV",
        "Scheduling: Daily, Weekly, Monthly, Quarterly",
        "Organization Branding: Logo, colors, custom footer"
    ], "Standardized reporting with reduced manual effort")

    add_content_slide(prs, "Platform Capabilities", [
        "Role-Based Access: Admin, Manager, Viewer roles",
        "Multi-Organization Support: Users can belong to multiple orgs",
        "Audit Logging: Track all user actions for compliance",
        "Responsive Design: Desktop, tablet, and mobile",
        "Dark/Light Mode: Customizable visual theme",
        "Real-Time Updates: Automatic data refresh"
    ])

    add_table_slide(prs, "Summary by User Role",
        ["Role", "Primary Modules"],
        [
            ["Executive", "Overview, AI Insights, Reports"],
            ["Procurement Manager", "Pareto, Stratification, Contracts, Maverick"],
            ["Category Manager", "Categories, Suppliers, Tail Spend"],
            ["AP/Finance", "Invoice Aging, 3-Way Matching, Supplier Payments"],
            ["Operations", "P2P Cycle, Requisitions, Purchase Orders"]
        ]
    )

    add_content_slide(prs, "Next Steps", [
        "1. Demo Session — Schedule hands-on walkthrough of key modules",
        "2. Data Requirements — Identify data sources for initial load",
        "3. User Provisioning — Define roles and access levels",
        "4. Report Templates — Customize report branding and schedules"
    ])

    add_title_slide(prs, "Thank You", "Questions?")

    prs.save("Versatex_Analytics_Management_Introduction.pptx")
    print("PowerPoint generated: Versatex_Analytics_Management_Introduction.pptx")


if __name__ == "__main__":
    build_pptx()
